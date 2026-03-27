"""
Document Ingestion Script
Processes documents from the documents folder and adds them to the vector store
"""

import sys
import os
from pathlib import Path

# Add project root to path
project_root = Path(__file__).parent.parent.parent.parent
sys.path.insert(0, str(project_root))

import logging
from typing import List, Dict, Any
from app.services.knowledge_base.vector_store_service import VectorStoreService
from app.services.llm_service import LLMService
import PyPDF2
import uuid
import fitz  # PyMuPDF for better PDF handling with images
from PIL import Image
import io
import pytesseract
import os
from pptx import Presentation

logging.basicConfig(
    level=logging.INFO,
    format='%(asctime)s - %(name)s - %(levelname)s - %(message)s'
)
logger = logging.getLogger(__name__)


class DocumentProcessor:
    """Process and ingest documents into vector store"""
    
    def __init__(self, use_ocr: bool = True):
        self.vector_store = VectorStoreService()
        self.llm_service = LLMService()
        self.use_ocr = use_ocr
        
        # Check if Tesseract is available
        if use_ocr:
            try:
                pytesseract.get_tesseract_version()
                logger.info("✅ OCR (Tesseract) available for image text extraction")
            except Exception:
                logger.warning("⚠️ Tesseract OCR not found. Images will be skipped.")
                logger.warning("   Install from: https://github.com/UB-Mannheim/tesseract/wiki")
                self.use_ocr = False
        
        logger.info("✅ Document Processor initialized")
    
    def extract_text_from_images(self, pdf_path: str) -> str:
        """Extract text from images in PDF using OCR"""
        if not self.use_ocr:
            return ""
        
        try:
            logger.info("🖼️ Extracting text from images using OCR...")
            ocr_text = ""
            image_count = 0
            
            # Open PDF with PyMuPDF (better image handling)
            pdf_document = fitz.open(pdf_path)
            
            for page_num in range(len(pdf_document)):
                page = pdf_document[page_num]
                
                # Get images from page
                image_list = page.get_images(full=True)
                
                for img_index, img in enumerate(image_list):
                    try:
                        xref = img[0]
                        base_image = pdf_document.extract_image(xref)
                        image_bytes = base_image["image"]
                        
                        # Convert to PIL Image
                        image = Image.open(io.BytesIO(image_bytes))
                        
                        # Perform OCR
                        text = pytesseract.image_to_string(image)
                        
                        if text.strip():
                            ocr_text += f"\n[Image {image_count + 1} from page {page_num + 1}]\n"
                            ocr_text += text.strip()
                            ocr_text += "\n\n"
                            image_count += 1
                    
                    except Exception as e:
                        logger.debug(f"   Could not process image {img_index} on page {page_num}: {e}")
                        continue
            
            pdf_document.close()
            
            if image_count > 0:
                logger.info(f"   ✅ Extracted text from {image_count} images ({len(ocr_text)} characters)")
            else:
                logger.info("   ℹ️ No text found in images")
            
            return ocr_text
        
        except Exception as e:
            logger.warning(f"⚠️ Error extracting text from images: {e}")
            return ""
    
    def extract_text_from_pdf(self, pdf_path: str) -> str:
        """Extract text content from PDF file (both text and images)"""
        try:
            logger.info(f"📄 Extracting text from: {pdf_path}")
            text = ""
            
            # Method 1: Extract regular text using PyPDF2
            logger.info("   Extracting regular text...")
            with open(pdf_path, 'rb') as file:
                pdf_reader = PyPDF2.PdfReader(file)
                num_pages = len(pdf_reader.pages)
                
                logger.info(f"   Found {num_pages} pages")
                
                for page_num in range(num_pages):
                    page = pdf_reader.pages[page_num]
                    page_text = page.extract_text()
                    if page_text:
                        text += page_text
                        text += "\n\n"
            
            text_chars = len(text)
            logger.info(f"   ✅ Extracted {text_chars} characters from regular text")
            
            # Method 2: Extract text from images using OCR
            if self.use_ocr:
                image_text = self.extract_text_from_images(pdf_path)
                if image_text:
                    text += "\n\n[TEXT FROM IMAGES]\n\n"
                    text += image_text
                    logger.info(f"   ✅ Total text: {len(text)} characters (including OCR)")
            
            return text.strip()
        
        except Exception as e:
            logger.error(f"❌ Error extracting PDF text: {e}")
            return ""
    
    def chunk_text(self, text: str, chunk_size: int = 1000, overlap: int = 200) -> List[str]:
        """Split text into overlapping chunks"""
        chunks = []
        start = 0
        
        while start < len(text):
            end = start + chunk_size
            chunk = text[start:end]
            
            # Try to break at sentence boundary
            if end < len(text):
                last_period = chunk.rfind('.')
                last_newline = chunk.rfind('\n')
                break_point = max(last_period, last_newline)
                
                if break_point > chunk_size * 0.7:  # At least 70% of chunk size
                    end = start + break_point + 1
                    chunk = text[start:end]
            
            chunks.append(chunk.strip())
            start = end - overlap
        
        return [c for c in chunks if len(c.strip()) > 50]  # Filter very small chunks
    
    def get_embedding(self, text: str) -> List[float]:
        """Get embedding vector for text using LLM service"""
        try:
            # Use the actual LLM service for embeddings
            embedding = self.llm_service.generate_embedding(text)
            return embedding
        except Exception as e:
            logger.warning(f"⚠️ Error generating embedding: {e}")
            logger.info("   Using fallback hash-based embedding")
            
            # Fallback to hash-based approach (1536 dimensions - matches OpenAI)
            import hashlib
            text_hash = hashlib.sha256(text.encode()).hexdigest()
            embedding = []
            
            for i in range(0, len(text_hash), 2):
                byte_val = int(text_hash[i:i+2], 16)
                embedding.append(float(byte_val) / 255.0)
            
            # Pad to 1536 dimensions (matches OpenAI text-embedding-3-small)
            while len(embedding) < 1536:
                embedding.append(0.0)
            
            return embedding[:1536]
    
    def extract_text_from_pptx(self, pptx_path: str) -> str:
        """Extract text from PowerPoint presentation"""
        try:
            logger.info(f"📊 Extracting text from PowerPoint: {Path(pptx_path).name}")
            text_content = []
            
            # Open presentation
            prs = Presentation(pptx_path)
            
            # Iterate through slides
            for slide_num, slide in enumerate(prs.slides, 1):
                slide_text = [f"\n--- Slide {slide_num} ---\n"]
                
                # Extract text from all shapes in the slide
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text:
                        slide_text.append(shape.text)
                    
                    # Extract text from tables
                    if shape.has_table:
                        table = shape.table
                        for row in table.rows:
                            row_text = []
                            for cell in row.cells:
                                if cell.text:
                                    row_text.append(cell.text)
                            if row_text:
                                slide_text.append(" | ".join(row_text))
                
                # Add slide content if it has text
                if len(slide_text) > 1:  # More than just the slide header
                    text_content.append("\n".join(slide_text))
            
            full_text = "\n\n".join(text_content)
            
            logger.info(f"   ✅ Extracted {len(full_text)} characters from {len(prs.slides)} slides")
            
            return full_text
        
        except Exception as e:
            logger.error(f"❌ Error extracting text from PowerPoint: {e}", exc_info=True)
            return ""
    
    def ingest_pptx(self, pptx_path: str, category: str = "documentation") -> Dict[str, Any]:
        """Ingest a PowerPoint document into vector store"""
        try:
            file_path = Path(pptx_path)
            logger.info(f"📥 Ingesting PowerPoint: {file_path.name}")
            
            # Extract text
            text = self.extract_text_from_pptx(pptx_path)
            
            if not text:
                logger.warning(f"⚠️ No text extracted from {file_path.name}")
                return {"status": "error", "message": "No text extracted"}
            
            # Split into chunks
            logger.info("✂️ Splitting into chunks...")
            chunks = self.chunk_text(text, chunk_size=1000, overlap=200)
            logger.info(f"   Created {len(chunks)} chunks")
            
            # Process each chunk
            documents = []
            for i, chunk in enumerate(chunks):
                doc_id = f"{file_path.stem}_chunk_{i}_{uuid.uuid4().hex[:8]}"
                
                # Get embedding
                embedding = self.get_embedding(chunk)
                
                # Create document
                doc = {
                    "id": doc_id,
                    "content": chunk,
                    "embedding": embedding,
                    "metadata": {
                        "filename": file_path.name,
                        "document_type": "pptx",
                        "category": category,
                        "chunk_index": i,
                        "total_chunks": len(chunks)
                    }
                }
                documents.append(doc)
            
            # Add to vector store
            logger.info("💾 Adding to vector store...")
            self.vector_store.add_documents_batch(documents)
            
            logger.info(f"✅ Successfully ingested {file_path.name} ({len(chunks)} chunks)")
            
            return {
                "status": "success",
                "filename": file_path.name,
                "chunks": len(chunks),
                "total_characters": len(text)
            }
        
        except Exception as e:
            logger.error(f"❌ Error ingesting PowerPoint: {e}", exc_info=True)
            return {"status": "error", "message": str(e)}
    
    def ingest_pdf(self, pdf_path: str, category: str = "documentation") -> Dict[str, Any]:
        """Ingest a PDF document into vector store"""
        try:
            file_path = Path(pdf_path)
            logger.info(f"📥 Ingesting document: {file_path.name}")
            
            # Extract text
            text = self.extract_text_from_pdf(pdf_path)
            
            if not text:
                logger.warning(f"⚠️ No text extracted from {file_path.name}")
                return {"status": "error", "message": "No text extracted"}
            
            # Split into chunks
            logger.info("✂️ Splitting into chunks...")
            chunks = self.chunk_text(text, chunk_size=1000, overlap=200)
            logger.info(f"   Created {len(chunks)} chunks")
            
            # Process each chunk
            documents = []
            for i, chunk in enumerate(chunks):
                doc_id = f"{file_path.stem}_chunk_{i}_{uuid.uuid4().hex[:8]}"
                
                # Get embedding
                embedding = self.get_embedding(chunk)
                
                # Create document
                doc = {
                    "id": doc_id,
                    "content": chunk,
                    "embedding": embedding,
                    "metadata": {
                        "filename": file_path.name,
                        "document_type": "pdf",
                        "category": category,
                        "chunk_index": i,
                        "total_chunks": len(chunks)
                    }
                }
                documents.append(doc)
            
            # Add to vector store
            logger.info("💾 Adding to vector store...")
            self.vector_store.add_documents_batch(documents)
            
            logger.info(f"✅ Successfully ingested {file_path.name} ({len(chunks)} chunks)")
            
            return {
                "status": "success",
                "filename": file_path.name,
                "chunks": len(chunks),
                "total_characters": len(text)
            }
        
        except Exception as e:
            logger.error(f"❌ Error ingesting document: {e}", exc_info=True)
            return {"status": "error", "message": str(e)}
    
    def ingest_folder(self, folder_path: str, category: str = "documentation"):
        """Ingest all PDF documents from a folder"""
        folder = Path(folder_path)
        
        if not folder.exists():
            logger.error(f"❌ Folder not found: {folder_path}")
            return
        
        logger.info(f"📁 Processing documents from: {folder_path}")
        logger.info("=" * 80)
        
        pdf_files = list(folder.glob("*.pdf"))
        
        if not pdf_files:
            logger.warning(f"⚠️ No PDF files found in {folder_path}")
            return
        
        logger.info(f"Found {len(pdf_files)} PDF files\n")
        
        results = []
        for pdf_file in pdf_files:
            result = self.ingest_pdf(str(pdf_file), category)
            results.append(result)
            print()  # Blank line between files
        
        # Summary
        logger.info("=" * 80)
        logger.info("📊 INGESTION SUMMARY")
        logger.info("=" * 80)
        
        successful = [r for r in results if r.get("status") == "success"]
        failed = [r for r in results if r.get("status") == "error"]
        
        logger.info(f"✅ Successful: {len(successful)}")
        logger.info(f"❌ Failed: {len(failed)}")
        
        if successful:
            total_chunks = sum(r.get("chunks", 0) for r in successful)
            logger.info(f"📦 Total chunks created: {total_chunks}")
        
        if failed:
            logger.info("\n❌ Failed documents:")
            for r in failed:
                logger.info(f"   - {r.get('filename', 'unknown')}: {r.get('message', 'unknown error')}")


def main():
    """Main function"""
    import argparse
    
    parser = argparse.ArgumentParser(description="Ingest documents into NEO knowledge base")
    parser.add_argument(
        "--folder",
        type=str,
        default="app/modules/neo_chatbot/data/documents",
        help="Folder containing documents to ingest"
    )
    parser.add_argument(
        "--category",
        type=str,
        default="documentation",
        help="Document category (documentation, code, proposal, support)"
    )
    parser.add_argument(
        "--file",
        type=str,
        help="Ingest a single file instead of folder"
    )
    parser.add_argument(
        "--no-ocr",
        action="store_true",
        help="Disable OCR for images (faster but skips image text)"
    )
    
    args = parser.parse_args()
    
    processor = DocumentProcessor(use_ocr=not args.no_ocr)
    
    print("\n" + "=" * 80)
    print("NEO DOCUMENT INGESTION")
    print("=" * 80 + "\n")
    
    if args.file:
        # Ingest single file
        result = processor.ingest_pdf(args.file, args.category)
        print(f"\nResult: {result}")
    else:
        # Ingest folder
        processor.ingest_folder(args.folder, args.category)
    
    print("\n" + "=" * 80)
    print("✅ INGESTION COMPLETE")
    print("=" * 80 + "\n")


if __name__ == "__main__":
    main()
