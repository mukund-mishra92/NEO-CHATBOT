"""
PPT Extractor — Multi-modal PowerPoint extraction.

Extracts:
  • Slide text (titles, body, text boxes)
  • Tables as markdown
  • Grouped shapes
  • Images with optional Vision API descriptions
  • Slide notes
"""

from __future__ import annotations

import io
import logging
from pathlib import Path
from typing import List, Optional

from ..models import ExtractedDocument, ExtractedElement, ExtractedPage, ImageReference

logger = logging.getLogger(__name__)


class PPTExtractor:
    """Extract structured content from .pptx files."""

    def __init__(self, *, image_describer=None, image_store=None):
        self.image_describer = image_describer
        self.image_store = image_store

    def extract(self, file_path: str) -> ExtractedDocument:
        """Extract all content from a PPTX file."""
        path = Path(file_path)
        doc = ExtractedDocument(
            source_path=str(path),
            file_type="pptx",
            title=path.stem,
            metadata={"filename": path.name},
        )

        try:
            from pptx import Presentation
        except ImportError:
            logger.error("python-pptx not installed — run: pip install python-pptx")
            return doc

        try:
            prs = Presentation(file_path)
            for slide_num, slide in enumerate(prs.slides, start=1):
                ep = ExtractedPage(page_number=slide_num)

                # ── Slide title ──
                slide_title = ""
                if slide.shapes.title:
                    slide_title = slide.shapes.title.text.strip()
                    if slide_title:
                        ep.elements.append(ExtractedElement(
                            content=slide_title,
                            element_type="heading",
                            page_number=slide_num,
                            metadata={"heading_level": 1, "slide_title": True},
                        ))

                # ── Extract notes early (needed for image context) ──
                slide_notes = ""
                if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
                    slide_notes = slide.notes_slide.notes_text_frame.text.strip()

                # ── Shapes ──
                for shape in slide.shapes:
                    self._process_shape(
                        shape, slide_num, ep, path.name,
                        slide_title=slide_title, slide_notes=slide_notes,
                    )

                # ── Slide notes element ──
                if slide_notes and len(slide_notes) > 5:
                    ep.elements.append(ExtractedElement(
                        content=f"[Speaker Notes: {slide_notes}]",
                        element_type="text",
                        page_number=slide_num,
                        metadata={"is_notes": True},
                    ))

                ep.composite_text = "\n\n".join(e.content for e in ep.elements)
                ep.raw_text = ep.composite_text
                doc.pages.append(ep)

            doc.metadata["total_pages"] = len(doc.pages)
            logger.info(
                f"✅ Extracted {len(doc.pages)} slides, "
                f"{sum(len(p.elements) for p in doc.pages)} elements from {path.name}"
            )

        except Exception as exc:
            logger.error(f"PPT extraction failed for {path.name}: {exc}", exc_info=True)

        return doc

    # ────────────────────────────────────────────────────
    #  Shape processing
    # ────────────────────────────────────────────────────
    # ── EMU → pixel conversion (python-pptx uses EMU; 914400 EMU = 1 inch) ──
    _EMU_PER_PIXEL = 914400 / 96  # 96 DPI

    def _process_shape(self, shape, slide_num: int, ep: ExtractedPage, filename: str,
                       *, slide_title: str = "", slide_notes: str = ""):
        """Process a single shape."""
        try:
            # Table
            if shape.has_table:
                md = self._table_to_markdown(shape.table)
                if md:
                    ep.elements.append(ExtractedElement(
                        content=md,
                        element_type="table",
                        page_number=slide_num,
                        metadata={
                            "rows": len(shape.table.rows),
                            "cols": len(shape.table.columns),
                        },
                    ))
                return

            # Text frame
            if shape.has_text_frame:
                text = shape.text_frame.text.strip()
                if text and len(text) > 2:
                    # Skip if it's the title (already captured)
                    el_type = "text"
                    # Check for bullet/list patterns
                    lines = text.split("\n")
                    if all(
                        line.strip().startswith(("•", "-", "●", "→", "▪", "✓", "✗"))
                        for line in lines
                        if line.strip()
                    ):
                        el_type = "list"

                    ep.elements.append(ExtractedElement(
                        content=text,
                        element_type=el_type,
                        page_number=slide_num,
                        metadata={"shape_name": shape.name},
                    ))
                return

            # Group shape — recurse
            if shape.shape_type is not None:
                try:
                    from pptx.enum.shapes import MSO_SHAPE_TYPE
                    if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                        for child in shape.shapes:
                            self._process_shape(
                                child, slide_num, ep, filename,
                                slide_title=slide_title, slide_notes=slide_notes,
                            )
                        return
                except Exception:
                    pass

            # Image / picture
            if hasattr(shape, "image") and shape.image:
                try:
                    img_bytes = shape.image.blob
                    if len(img_bytes) < 2000:
                        return  # skip tiny images

                    # Save to disk
                    image_path = ""
                    if self.image_store:
                        img_idx = len([
                            e for e in ep.elements if e.element_type == "image"
                        ])
                        image_path = self.image_store.save(
                            img_bytes, filename, slide_num, img_idx,
                        )

                    # Vision API description
                    desc = ""
                    if self.image_describer:
                        desc = self.image_describer.describe(
                            img_bytes,
                            context=f"Slide {slide_num} of {filename}"
                        )

                    # ── Build contextual caption from slide context ──
                    context_parts = []
                    if slide_title:
                        context_parts.append(slide_title)
                    # Collect text already extracted on this slide
                    slide_texts = [
                        e.content for e in ep.elements
                        if e.element_type in ("text", "list") and len(e.content) > 5
                    ]
                    if slide_texts:
                        context_parts.append(" ".join(slide_texts)[:300])
                    if slide_notes:
                        context_parts.append(slide_notes[:200])
                    if not context_parts:
                        context_parts.append(f"Slide {slide_num} of {filename}")
                    caption = " | ".join(context_parts)

                    # ── Convert EMU to pixels ──
                    raw_w = int(getattr(shape, "width", 0) or 0)
                    raw_h = int(getattr(shape, "height", 0) or 0)
                    px_w = int(raw_w / self._EMU_PER_PIXEL) if raw_w > 20000 else raw_w
                    px_h = int(raw_h / self._EMU_PER_PIXEL) if raw_h > 20000 else raw_h

                    # Build ImageReference
                    img_ref = ImageReference(
                        image_path=image_path,
                        page_number=slide_num,
                        caption=caption,
                        description=desc,
                        width=px_w,
                        height=px_h,
                    )

                    display_text = desc or caption
                    ep.elements.append(ExtractedElement(
                        content=f"[Image: {display_text}]",
                        element_type="image",
                        page_number=slide_num,
                        metadata={
                            "description": desc,
                            "caption": caption,
                            "shape_name": shape.name,
                            "image_path": image_path,
                            "image_ref": img_ref.to_dict(),
                        },
                    ))
                except Exception:
                    pass

        except Exception as exc:
            logger.debug(f"Error processing shape on slide {slide_num}: {exc}")

    # ────────────────────────────────────────────────────
    #  Table → Markdown
    # ────────────────────────────────────────────────────
    @staticmethod
    def _table_to_markdown(table) -> str:
        """Convert a python-pptx Table to Markdown."""
        rows_data = []
        for row in table.rows:
            cells = [cell.text.replace("\n", " ").strip() for cell in row.cells]
            rows_data.append(cells)
        if not rows_data:
            return ""

        header = rows_data[0]
        md = "| " + " | ".join(header) + " |\n"
        md += "| " + " | ".join("---" for _ in header) + " |\n"
        for row in rows_data[1:]:
            padded = row + [""] * (len(header) - len(row)) if len(row) < len(header) else row[:len(header)]
            md += "| " + " | ".join(padded) + " |\n"
        return md.strip()
